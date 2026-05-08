from pptx import Presentation
from pptx.util import Cm, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches

# 색상 정의
RED = RGBColor(180, 30, 30)
DARK = RGBColor(30, 30, 30)
GRAY_BG = RGBColor(248, 248, 248)
WHITE = RGBColor(255, 255, 255)
DARK_CARD = RGBColor(50, 50, 50)
GRAY_TEXT = RGBColor(160, 160, 160)

FONT_NAME = '맑은 고딕'

W = Cm(21)
H = Cm(29.7)

prs = Presentation()
prs.slide_width = W
prs.slide_height = H

blank_layout = prs.slide_layouts[6]  # blank

def set_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_text_box(slide, text, x, y, w, h,
                 font_size=14, bold=False, italic=False,
                 color=WHITE, align=PP_ALIGN.LEFT,
                 line_spacing=None, wrap=True):
    txBox = slide.shapes.add_textbox(x, y, w, h)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    if line_spacing:
        from pptx.util import Pt as PT
        from pptx.oxml.ns import qn
        from lxml import etree
        pPr = p._pPr if p._pPr is not None else p._p.get_or_add_pPr()
        lnSpc = etree.SubElement(pPr, qn('a:lnSpc'))
        spcPts = etree.SubElement(lnSpc, qn('a:spcPts'))
        spcPts.set('val', str(int(line_spacing * 100)))
    run = p.add_run()
    run.text = text
    run.font.name = FONT_NAME
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox

def add_rect(slide, x, y, w, h, fill_color, line_color=None):
    shape = slide.shapes.add_shape(
        1, x, y, w, h  # MSO_SHAPE_TYPE.RECTANGLE = 1
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
    else:
        shape.line.fill.background()
    return shape

def add_multiline_textbox(slide, lines, x, y, w, h,
                           font_sizes, bolds, colors,
                           align=PP_ALIGN.LEFT, line_spacing_pt=None):
    from pptx.oxml.ns import qn
    from lxml import etree
    txBox = slide.shapes.add_textbox(x, y, w, h)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, (text, fsize, fbold, fcolor) in enumerate(zip(lines, font_sizes, bolds, colors)):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = align
        if line_spacing_pt:
            pPr = p._p.get_or_add_pPr()
            lnSpc = etree.SubElement(pPr, qn('a:lnSpc'))
            spcPts = etree.SubElement(lnSpc, qn('a:spcPts'))
            spcPts.set('val', str(int(line_spacing_pt * 100)))
        run = p.add_run()
        run.text = text
        run.font.name = FONT_NAME
        run.font.size = Pt(fsize)
        run.font.bold = fbold
        run.font.color.rgb = fcolor
    return txBox

# ─────────────────────────────────────────────
# 슬라이드 1: 표지
# ─────────────────────────────────────────────
slide1 = prs.slides.add_slide(blank_layout)
set_bg(slide1, RED)

# 로고
add_text_box(slide1, '대기만성',
             Cm(0), Cm(3), W, Cm(1.5),
             font_size=28, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# 메인 카피
add_multiline_textbox(slide1,
    ['지금 중식 하시는데', '잘 안 되고 계신가요?'],
    Cm(1), Cm(8), Cm(19), Cm(5),
    [32, 32], [True, True], [WHITE, WHITE],
    align=PP_ALIGN.CENTER, line_spacing_pt=48)

# 서브
add_text_box(slide1,
    '가맹비 · 교육비 800만원 전액 면제로 전환하세요',
    Cm(1), Cm(15), Cm(19), Cm(1.5),
    font_size=16, color=WHITE, align=PP_ALIGN.CENTER)

# 구분선
add_rect(slide1, Cm(1), Cm(27), Cm(19), Cm(0.05), WHITE)

# 연락처
add_text_box(slide1, '대기만성   031-684-0095',
             Cm(1), Cm(27.3), Cm(19), Cm(1.2),
             font_size=14, color=WHITE, align=PP_ALIGN.CENTER)


# ─────────────────────────────────────────────
# 슬라이드 2: 공감 페이지
# ─────────────────────────────────────────────
slide2 = prs.slides.add_slide(blank_layout)
set_bg(slide2, WHITE)

# 상단 레드 바
add_rect(slide2, Cm(0), Cm(0), W, Cm(1.8), RED)
add_text_box(slide2, '지금 이런 상황 아니신가요?',
             Cm(0.5), Cm(0.2), Cm(20), Cm(1.4),
             font_size=16, bold=True, color=WHITE, align=PP_ALIGN.LEFT)

# 체크리스트
items = [
    '✓  주방장 월급만 500~600만원, 매출은 그만큼 안 나온다',
    '✓  주방장 눈치를 내가 본다',
    '✓  주방장이 나가면 매장이 멈춘다',
    '✓  장사는 되는데 손에 남는 게 없다',
]
y_start = Cm(3.5)
for i, item in enumerate(items):
    add_text_box(slide2, item,
                 Cm(1.5), y_start + i * Cm(4.5), Cm(18), Cm(2.5),
                 font_size=18, bold=False, color=DARK, align=PP_ALIGN.LEFT)
    if i < len(items) - 1:
        add_rect(slide2, Cm(1.5), y_start + (i+1)*Cm(4.5) - Cm(0.3),
                 Cm(18), Cm(0.04), RGBColor(220, 220, 220))

# 하단 레드 박스
add_rect(slide2, Cm(0), Cm(22.5), W, Cm(2.5), RED)
add_text_box(slide2, '기술자에게 끌려다니는 구조, 바꿀 수 있습니다.',
             Cm(0.5), Cm(22.8), Cm(20), Cm(2),
             font_size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)


# ─────────────────────────────────────────────
# 슬라이드 3: 해결책 페이지
# ─────────────────────────────────────────────
slide3 = prs.slides.add_slide(blank_layout)
set_bg(slide3, GRAY_BG)

# 상단 레드 바
add_rect(slide3, Cm(0), Cm(0), W, Cm(1.8), RED)
add_text_box(slide3, '대기만성은 처음부터 다르게 설계됐습니다',
             Cm(0.5), Cm(0.2), Cm(20), Cm(1.4),
             font_size=16, bold=True, color=WHITE)

# 3개 박스
box_data = [
    ('🚚 물류 안정화', '표준화된 식자재\n매일 동일 공급'),
    ('📋 3일 완성 교육', '요리 경험 불필요\n누가 만들어도 같은 맛'),
    ('📦 35% 원가 설계', '표준화된 수익 구조\n안정적 마진'),
]
bx = Cm(0.7)
bw = Cm(6)
by = Cm(5)
bh = Cm(8)
gap = Cm(0.4)

for i, (title, content) in enumerate(box_data):
    x = bx + i * (bw + gap)
    # 흰색 박스
    add_rect(slide3, x, by, bw, bh, WHITE)
    # 레드 상단 포인트
    add_rect(slide3, x, by, bw, Cm(0.35), RED)
    # 제목
    add_text_box(slide3, title,
                 x + Cm(0.2), by + Cm(0.6), bw - Cm(0.4), Cm(1.5),
                 font_size=15, bold=True, color=RED, align=PP_ALIGN.CENTER)
    # 내용
    add_text_box(slide3, content,
                 x + Cm(0.2), by + Cm(2.3), bw - Cm(0.4), Cm(4),
                 font_size=13, color=DARK, align=PP_ALIGN.CENTER)

# 하단 강조
add_text_box(slide3, '직영 13개점을 직접 운영하며 검증한 시스템입니다',
             Cm(1), Cm(22), Cm(19), Cm(1.5),
             font_size=15, italic=True, color=DARK, align=PP_ALIGN.CENTER)


# ─────────────────────────────────────────────
# 슬라이드 4: 실제 데이터 페이지
# ─────────────────────────────────────────────
slide4 = prs.slides.add_slide(blank_layout)
set_bg(slide4, DARK)

# 제목
add_text_box(slide4, '실제 직영점 운영 데이터입니다',
             Cm(1), Cm(2), Cm(19), Cm(1.5),
             font_size=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text_box(slide4, '조작 없는 실제 수치입니다',
             Cm(1), Cm(3.8), Cm(19), Cm(1),
             font_size=13, color=GRAY_TEXT, align=PP_ALIGN.CENTER)

# 3개 카드
cards = [
    ('고덕점 📍 평택', '월 평균', '1억 1,000만원'),
    ('공도점 📍 안성', '월 평균', '7,900만원'),
    ('청당점 📍 천안', '월 평균', '4,400만원'),
]
cx = Cm(0.7)
cw = Cm(6)
cy = Cm(6.5)
ch = Cm(9)

for i, (name, label, amount) in enumerate(cards):
    x = cx + i * (cw + Cm(0.4))
    add_rect(slide4, x, cy, cw, ch, DARK_CARD)
    add_rect(slide4, x, cy, cw, Cm(0.35), RED)
    add_text_box(slide4, name,
                 x + Cm(0.2), cy + Cm(0.7), cw - Cm(0.4), Cm(1.2),
                 font_size=14, bold=False, color=WHITE, align=PP_ALIGN.CENTER)
    add_text_box(slide4, label,
                 x + Cm(0.2), cy + Cm(2.2), cw - Cm(0.4), Cm(1),
                 font_size=13, color=GRAY_TEXT, align=PP_ALIGN.CENTER)
    add_text_box(slide4, amount,
                 x + Cm(0.1), cy + Cm(3.5), cw - Cm(0.2), Cm(2),
                 font_size=20, bold=True, color=RED, align=PP_ALIGN.CENTER)

# 구분선
add_rect(slide4, Cm(1), Cm(24), Cm(19), Cm(0.05), GRAY_TEXT)

# 하단 리뷰
add_text_box(slide4,
    '배달 리뷰 9,999건+   |   전 플랫폼 평균 별점 4.9점',
    Cm(1), Cm(24.5), Cm(19), Cm(1.5),
    font_size=15, color=WHITE, align=PP_ALIGN.CENTER)


# ─────────────────────────────────────────────
# 슬라이드 5: 비교표 페이지
# ─────────────────────────────────────────────
slide5 = prs.slides.add_slide(blank_layout)
set_bg(slide5, WHITE)

# 상단 레드 바
add_rect(slide5, Cm(0), Cm(0), W, Cm(1.8), RED)
add_text_box(slide5, '일반 중식당과 무엇이 다른가요?',
             Cm(0.5), Cm(0.2), Cm(20), Cm(1.4),
             font_size=16, bold=True, color=WHITE)

# 표 설정
table_x = Cm(1)
table_y = Cm(3)
table_w = Cm(19)
table_h = Cm(18)
col_w = [Cm(5), Cm(7), Cm(7)]
row_h = Cm(3.2)

headers = ['항목', '❌ 일반 중식당', '✅ 대기만성']
rows = [
    ('주방장 인건비', '월 500~600만원 필수', '불필요'),
    ('맛 일관성', '사람에 따라 달라짐', '시스템이 유지'),
    ('직원 퇴사 리스크', '맛 붕괴 위험', '영향 없음'),
    ('신규직원 교육', '수개월 소요', '단 3일 완성'),
]

# 헤더 행
x_cursor = table_x
for j, (header, cw) in enumerate(zip(headers, col_w)):
    add_rect(slide5, x_cursor, table_y, cw, row_h, DARK)
    add_text_box(slide5, header,
                 x_cursor + Cm(0.2), table_y + Cm(0.8), cw - Cm(0.4), Cm(1.5),
                 font_size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    x_cursor += cw

# 데이터 행
row_bgs = [WHITE, GRAY_BG]
for i, row in enumerate(rows):
    y = table_y + (i + 1) * row_h
    x_cursor = table_x
    bg = row_bgs[i % 2]
    for j, (cell, cw) in enumerate(zip(row, col_w)):
        add_rect(slide5, x_cursor, y, cw, row_h, bg,
                 line_color=RGBColor(220, 220, 220))
        is_dgms = (j == 2)
        add_text_box(slide5, cell,
                     x_cursor + Cm(0.2), y + Cm(0.7), cw - Cm(0.4), Cm(1.8),
                     font_size=14, bold=is_dgms,
                     color=RED if is_dgms else DARK,
                     align=PP_ALIGN.CENTER)
        x_cursor += cw


# ─────────────────────────────────────────────
# 슬라이드 6: CTA 페이지
# ─────────────────────────────────────────────
slide6 = prs.slides.add_slide(blank_layout)
set_bg(slide6, RED)

# 상단
add_text_box(slide6, '기존 중식집 전환 사장님께만',
             Cm(1), Cm(2.5), Cm(19), Cm(1.5),
             font_size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# 중앙 큰 텍스트
add_text_box(slide6, '선착순 30호점 한정 혜택',
             Cm(1), Cm(4.5), Cm(19), Cm(2),
             font_size=28, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# 혜택 박스
box_y = Cm(8)
box_h = Cm(8)
add_rect(slide6, Cm(2), box_y, Cm(17), box_h, WHITE)

benefit_lines = [
    '가맹비  500만원  →  전액 면제',
    '교육비  300만원  →  전액 면제',
    '─────────────────',
    '합계   800만원  →  0원',
]
for i, line in enumerate(benefit_lines):
    add_text_box(slide6, line,
                 Cm(2.5), box_y + Cm(0.8) + i * Cm(1.8), Cm(16), Cm(1.6),
                 font_size=18, bold=True, color=RED, align=PP_ALIGN.CENTER)

# 하단 텍스트
add_text_box(slide6, '부담 없이 먼저 상담받아 보세요',
             Cm(1), Cm(18), Cm(19), Cm(1.5),
             font_size=16, color=WHITE, align=PP_ALIGN.CENTER)

# 구분선
add_rect(slide6, Cm(1), Cm(26.5), Cm(19), Cm(0.05), WHITE)

# 연락처
add_text_box(slide6, '대기만성   031-684-0095   |   chc9443@naver.com',
             Cm(1), Cm(26.8), Cm(19), Cm(1),
             font_size=13, color=WHITE, align=PP_ALIGN.CENTER)
add_text_box(slide6, '주식회사 대기만성홀딩스',
             Cm(1), Cm(27.9), Cm(19), Cm(1),
             font_size=13, color=WHITE, align=PP_ALIGN.CENTER)

# ─────────────────────────────────────────────
# 저장
# ─────────────────────────────────────────────
prs.save('/Users/user/Desktop/DGMS_homepage/dgms_brochure.pptx')
print("완료. 총 6슬라이드 생성됨")
